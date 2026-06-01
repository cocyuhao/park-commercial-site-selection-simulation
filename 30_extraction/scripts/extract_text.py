from __future__ import annotations

import json
from datetime import datetime
from pathlib import Path


ROOT = Path(__file__).resolve().parents[2]
RAW_DIR = ROOT / "20_raw_data"
PDF_TEXT_DIR = ROOT / "30_extraction" / "pdf_text"
PPT_TEXT_DIR = ROOT / "30_extraction" / "ppt_text"
LOG_PATH = ROOT / "30_extraction" / "extraction_logs" / "text_extraction_log.json"


def safe_stem(path: Path) -> str:
    return path.stem.replace("/", "_").replace("\\", "_")


def extract_pdf(path: Path) -> dict:
    import fitz

    doc = fitz.open(str(path))
    pages = []
    for idx, page in enumerate(doc, start=1):
        text = page.get_text("text") or ""
        pages.append({"page": idx, "text_length": len(text.strip()), "text": text})

    out_path = PDF_TEXT_DIR / f"{safe_stem(path)}.json"
    out_path.write_text(json.dumps({
        "source_file": path.relative_to(ROOT).as_posix(),
        "page_count": len(pages),
        "pages": pages,
    }, ensure_ascii=False, indent=2), encoding="utf-8")
    return {"file": path.name, "type": "pdf", "status": "ok", "page_count": len(pages), "output": out_path.relative_to(ROOT).as_posix()}


def extract_pptx(path: Path) -> dict:
    from pptx import Presentation

    prs = Presentation(str(path))
    slides = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        table_count = 0
        picture_count = 0
        chart_count = 0
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
            if getattr(shape, "has_table", False):
                table_count += 1
            if getattr(shape, "has_chart", False):
                chart_count += 1
            if getattr(shape, "shape_type", None) == 13:
                picture_count += 1
        joined = "\n".join(texts)
        slides.append({
            "slide": idx,
            "text_length": len(joined.strip()),
            "table_count": table_count,
            "picture_count": picture_count,
            "chart_count": chart_count,
            "text": joined,
        })

    out_path = PPT_TEXT_DIR / f"{safe_stem(path)}.json"
    out_path.write_text(json.dumps({
        "source_file": path.relative_to(ROOT).as_posix(),
        "slide_count": len(slides),
        "slides": slides,
    }, ensure_ascii=False, indent=2), encoding="utf-8")
    return {"file": path.name, "type": "pptx", "status": "ok", "slide_count": len(slides), "output": out_path.relative_to(ROOT).as_posix()}


def main() -> None:
    PDF_TEXT_DIR.mkdir(parents=True, exist_ok=True)
    PPT_TEXT_DIR.mkdir(parents=True, exist_ok=True)
    LOG_PATH.parent.mkdir(parents=True, exist_ok=True)

    log = {"started_at": datetime.now().isoformat(timespec="seconds"), "results": []}
    for path in sorted(RAW_DIR.rglob("*")):
        if not path.is_file() or path.name.startswith("~$"):
            continue
        suffix = path.suffix.lower()
        try:
            if suffix == ".pdf":
                log["results"].append(extract_pdf(path))
            elif suffix in {".pptx", ".ppt"}:
                log["results"].append(extract_pptx(path))
        except Exception as exc:
            log["results"].append({"file": path.name, "type": suffix, "status": "error", "error": repr(exc)})

    log["finished_at"] = datetime.now().isoformat(timespec="seconds")
    LOG_PATH.write_text(json.dumps(log, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"extracted {len(log['results'])} files -> {LOG_PATH}")


if __name__ == "__main__":
    main()
