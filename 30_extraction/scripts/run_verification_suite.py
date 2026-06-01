from __future__ import annotations

import csv
import json
import re
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Any

import fitz
from pptx import Presentation


ROOT = Path(__file__).resolve().parents[2]
RAW_DIR = ROOT / "20_raw_data"
PDF_TEXT_DIR = ROOT / "30_extraction" / "pdf_text"
PPT_TEXT_DIR = ROOT / "30_extraction" / "ppt_text"
QUALITY_DIR = ROOT / "40_quality_evidence"
VERIFY_DIR = QUALITY_DIR / "verification"

CATALOG_PATH = QUALITY_DIR / "data_catalog.csv"
PROFILE_PATH = QUALITY_DIR / "source_profile.csv"
KEYWORD_HITS_PATH = ROOT / "30_extraction" / "tables" / "keyword_hits.csv"
LEDGER_PATH = QUALITY_DIR / "evidence_ledger.csv"

SUPPORTED_SUFFIXES = {".pdf", ".ppt", ".pptx", ".csv", ".xlsx", ".xls", ".png", ".jpg", ".jpeg", ".webp"}
KEYWORDS = [
    "客流",
    "到访",
    "TGI",
    "POI",
    "热门到访",
    "品牌",
    "咖啡",
    "餐饮",
    "小吃",
    "快餐",
    "供需",
    "缺口",
    "外溢",
    "收入",
    "收益",
    "回收期",
    "转化",
    "消费",
]


@dataclass
class Check:
    check_id: str
    scope: str
    method: str
    status: str
    detail: str
    evidence_path: str = ""


def read_csv_rows(path: Path) -> list[dict[str, str]]:
    with path.open(encoding="utf-8-sig", newline="") as f:
        return list(csv.DictReader(f))


def write_csv(path: Path, rows: list[dict[str, Any]], fieldnames: list[str]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", encoding="utf-8-sig", newline="") as f:
        writer = csv.DictWriter(f, fieldnames=fieldnames)
        writer.writeheader()
        writer.writerows(rows)


def normalize(text: str) -> str:
    return re.sub(r"\s+", " ", text).strip()


def count_keywords(text: str) -> dict[str, int]:
    return {kw: text.count(kw) for kw in KEYWORDS}


def numeric_profile(text: str) -> dict[str, int]:
    clean = normalize(text)
    return {
        "number_count": len(re.findall(r"(?<![A-Za-z])[-+]?\d+(?:,\d{3})*(?:\.\d+)?", clean)),
        "percent_count": len(re.findall(r"[-+]?\d+(?:\.\d+)?%", clean)),
        "currency_count": len(re.findall(r"(?:￥|¥|元)\s*[-+]?\d+(?:,\d{3})*(?:\.\d+)?|[-+]?\d+(?:,\d{3})*(?:\.\d+)?\s*(?:元|万元|亿元)", clean)),
        "date_count": len(re.findall(r"20\d{6}|20\d{2}[-./年]\d{1,2}", clean)),
        "zero_like_count": len(re.findall(r"(?<!\d)0(?:\.0+)?(?!\d)", clean)),
    }


def table_like_score(text: str) -> int:
    clean = normalize(text)
    labels = ["排序", "排名", "占比", "指标", "数值", "TGI", "POI名称", "品牌名称", "扩展信息", "业态", "回收期"]
    score = sum(2 for label in labels if label in clean)
    score += min(10, len(re.findall(r"\d+(?:\.\d+)?%", clean)))
    score += min(10, len(re.findall(r"(?<![A-Za-z])\d{2,}(?![A-Za-z])", clean)))
    return score


def source_files() -> list[Path]:
    return sorted(
        [
            p
            for p in RAW_DIR.rglob("*")
            if p.is_file() and not p.name.startswith("~$") and p.suffix.lower() in SUPPORTED_SUFFIXES
        ]
    )


def check_catalog(checks: list[Check]) -> None:
    rows = read_csv_rows(CATALOG_PATH)
    raw_files = source_files()
    raw_rel = {p.relative_to(ROOT).as_posix(): p for p in raw_files}
    catalog_rel = {r["file_path"]: r for r in rows}

    status = "PASS" if set(raw_rel) == set(catalog_rel) else "FAIL"
    detail = f"raw_supported_files={len(raw_rel)}, catalog_rows={len(catalog_rel)}"
    if status == "FAIL":
        missing = sorted(set(raw_rel) - set(catalog_rel))
        extra = sorted(set(catalog_rel) - set(raw_rel))
        detail += f", missing_in_catalog={missing}, extra_in_catalog={extra}"
    checks.append(Check("CAT-001", "data_catalog", "raw file list equals catalog", status, detail, "40_quality_evidence/data_catalog.csv"))

    for idx, row in enumerate(rows, start=1):
        rel = row["file_path"]
        path = ROOT / rel
        exists = path.exists()
        size_ok = exists and str(path.stat().st_size) == str(row["size_bytes"])
        checks.append(
            Check(
                f"CAT-{idx+1:03d}",
                rel,
                "catalog path and size check",
                "PASS" if exists and size_ok else "FAIL",
                f"exists={exists}, size_ok={size_ok}",
                rel,
            )
        )


def verify_pdf(path: Path, checks: list[Check], numeric_rows: list[dict[str, Any]], table_rows: list[dict[str, Any]]) -> None:
    json_path = PDF_TEXT_DIR / f"{path.stem}.json"
    rel = path.relative_to(ROOT).as_posix()
    doc = fitz.open(str(path))
    data = json.loads(json_path.read_text(encoding="utf-8")) if json_path.exists() else {}
    pages = data.get("pages", [])

    checks.append(
        Check(
            f"PDF-{len(checks)+1:03d}",
            rel,
            "raw PDF page count vs extracted JSON",
            "PASS" if len(doc) == len(pages) else "FAIL",
            f"raw_pages={len(doc)}, json_pages={len(pages)}",
            json_path.relative_to(ROOT).as_posix() if json_path.exists() else "",
        )
    )

    raw_lengths = []
    json_lengths = []
    table_count = 0
    for i, page in enumerate(doc, start=1):
        raw_text = page.get_text("text") or ""
        json_text = pages[i - 1].get("text", "") if i - 1 < len(pages) else ""
        raw_lengths.append(len(raw_text.strip()))
        json_lengths.append(len(json_text.strip()))

        profile = numeric_profile(json_text)
        numeric_rows.append(
            {
                "source_file": rel,
                "unit_type": "page",
                "unit_index": i,
                "text_length": len(json_text.strip()),
                "table_like_score": table_like_score(json_text),
                **profile,
            }
        )

        try:
            tables = page.find_tables()
            page_table_count = len(tables.tables)
        except Exception:
            page_table_count = -1
        if page_table_count > 0:
            table_count += page_table_count
        if page_table_count != 0 or table_like_score(json_text) >= 12:
            table_rows.append(
                {
                    "source_file": rel,
                    "unit_type": "page",
                    "unit_index": i,
                    "pymupdf_table_count": page_table_count,
                    "table_like_score": table_like_score(json_text),
                    "numeric_count": profile["number_count"],
                    "percent_count": profile["percent_count"],
                    "currency_count": profile["currency_count"],
                    "snippet": normalize(json_text)[:220],
                }
            )

    checks.append(
        Check(
            f"PDF-{len(checks)+1:03d}",
            rel,
            "all JSON page text nonempty",
            "PASS" if all(x > 0 for x in json_lengths) else "WARN",
            f"nonempty_pages={sum(1 for x in json_lengths if x > 0)}/{len(json_lengths)}, total_text={sum(json_lengths)}",
            json_path.relative_to(ROOT).as_posix(),
        )
    )
    checks.append(
        Check(
            f"PDF-{len(checks)+1:03d}",
            rel,
            "raw extraction length equals saved JSON length",
            "PASS" if raw_lengths == json_lengths else "FAIL",
            f"raw_total={sum(raw_lengths)}, json_total={sum(json_lengths)}",
            json_path.relative_to(ROOT).as_posix(),
        )
    )
    checks.append(
        Check(
            f"PDF-{len(checks)+1:03d}",
            rel,
            "PyMuPDF table detection",
            "PASS" if table_count >= 0 else "WARN",
            f"detected_native_tables={table_count}, heuristic_candidates_recorded_in=table_candidates.csv",
            "40_quality_evidence/verification/table_candidates.csv",
        )
    )


def verify_ppt(path: Path, checks: list[Check], numeric_rows: list[dict[str, Any]], table_rows: list[dict[str, Any]]) -> None:
    json_path = PPT_TEXT_DIR / f"{path.stem}.json"
    rel = path.relative_to(ROOT).as_posix()
    prs = Presentation(str(path))
    data = json.loads(json_path.read_text(encoding="utf-8")) if json_path.exists() else {}
    slides = data.get("slides", [])

    checks.append(
        Check(
            f"PPT-{len(checks)+1:03d}",
            rel,
            "raw PPT slide count vs extracted JSON",
            "PASS" if len(prs.slides) == len(slides) else "FAIL",
            f"raw_slides={len(prs.slides)}, json_slides={len(slides)}",
            json_path.relative_to(ROOT).as_posix() if json_path.exists() else "",
        )
    )

    raw_text_lengths = []
    json_text_lengths = []
    raw_table_count = 0
    raw_picture_count = 0
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        table_count = 0
        picture_count = 0
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
            if getattr(shape, "has_table", False):
                table_count += 1
            if getattr(shape, "shape_type", None) == 13:
                picture_count += 1
        raw_table_count += table_count
        raw_picture_count += picture_count
        raw_text = "\n".join(texts)
        json_text = slides[i - 1].get("text", "") if i - 1 < len(slides) else ""
        raw_text_lengths.append(len(raw_text.strip()))
        json_text_lengths.append(len(json_text.strip()))

        profile = numeric_profile(json_text)
        numeric_rows.append(
            {
                "source_file": rel,
                "unit_type": "slide",
                "unit_index": i,
                "text_length": len(json_text.strip()),
                "table_like_score": table_like_score(json_text),
                **profile,
            }
        )
        if table_count > 0 or table_like_score(json_text) >= 12:
            table_rows.append(
                {
                    "source_file": rel,
                    "unit_type": "slide",
                    "unit_index": i,
                    "pymupdf_table_count": "",
                    "table_like_score": table_like_score(json_text),
                    "numeric_count": profile["number_count"],
                    "percent_count": profile["percent_count"],
                    "currency_count": profile["currency_count"],
                    "snippet": normalize(json_text)[:220],
                }
            )

    checks.append(
        Check(
            f"PPT-{len(checks)+1:03d}",
            rel,
            "all JSON slide text nonempty",
            "PASS" if all(x > 0 for x in json_text_lengths) else "WARN",
            f"nonempty_slides={sum(1 for x in json_text_lengths if x > 0)}/{len(json_text_lengths)}, total_text={sum(json_text_lengths)}",
            json_path.relative_to(ROOT).as_posix(),
        )
    )
    checks.append(
        Check(
            f"PPT-{len(checks)+1:03d}",
            rel,
            "raw PPT text length equals saved JSON length",
            "PASS" if raw_text_lengths == json_text_lengths else "FAIL",
            f"raw_total={sum(raw_text_lengths)}, json_total={sum(json_text_lengths)}",
            json_path.relative_to(ROOT).as_posix(),
        )
    )
    checks.append(
        Check(
            f"PPT-{len(checks)+1:03d}",
            rel,
            "native PPT tables and pictures",
            "PASS",
            f"native_tables={raw_table_count}, pictures={raw_picture_count}; table-like text recorded separately",
            "40_quality_evidence/verification/table_candidates.csv",
        )
    )


def verify_keyword_hits(checks: list[Check]) -> None:
    rows = read_csv_rows(KEYWORD_HITS_PATH)
    csv_counts: dict[str, int] = {}
    for row in rows:
        csv_counts[row["keyword"]] = csv_counts.get(row["keyword"], 0) + 1

    direct_counts = {kw: 0 for kw in KEYWORDS}
    for path in sorted(PDF_TEXT_DIR.glob("*.json")):
        data = json.loads(path.read_text(encoding="utf-8"))
        for page in data.get("pages", []):
            for kw, n in count_keywords(page.get("text", "")).items():
                direct_counts[kw] += n
    for path in sorted(PPT_TEXT_DIR.glob("*.json")):
        data = json.loads(path.read_text(encoding="utf-8"))
        for slide in data.get("slides", []):
            for kw, n in count_keywords(slide.get("text", "")).items():
                direct_counts[kw] += n

    # keyword_hits records up to 3 snippets per keyword per page/slide, so it may be lower than true occurrence counts.
    summary_rows = []
    for kw in KEYWORDS:
        summary_rows.append({"keyword": kw, "keyword_hits_rows": csv_counts.get(kw, 0), "direct_occurrences": direct_counts.get(kw, 0)})

    out = VERIFY_DIR / "keyword_recount.csv"
    write_csv(out, summary_rows, ["keyword", "keyword_hits_rows", "direct_occurrences"])
    checks.append(
        Check(
            f"KEY-{len(checks)+1:03d}",
            "keyword_hits",
            "recompute keyword occurrence counts from extracted JSON",
            "PASS" if len(rows) > 0 and sum(direct_counts.values()) >= len(rows) else "WARN",
            f"keyword_hit_rows={len(rows)}, direct_occurrences={sum(direct_counts.values())}; hit rows are capped snippets per unit",
            out.relative_to(ROOT).as_posix(),
        )
    )


def verify_profile(checks: list[Check]) -> None:
    rows = read_csv_rows(PROFILE_PATH)
    expected_sources = set()
    for path in sorted(PDF_TEXT_DIR.glob("*.json")):
        expected_sources.add(json.loads(path.read_text(encoding="utf-8")).get("source_file"))
    for path in sorted(PPT_TEXT_DIR.glob("*.json")):
        expected_sources.add(json.loads(path.read_text(encoding="utf-8")).get("source_file"))
    profile_sources = {r["source_file"] for r in rows}
    checks.append(
        Check(
            f"PRO-{len(checks)+1:03d}",
            "source_profile",
            "profile rows match extracted JSON sources",
            "PASS" if expected_sources == profile_sources else "FAIL",
            f"profile_rows={len(rows)}, extracted_sources={len(expected_sources)}",
            PROFILE_PATH.relative_to(ROOT).as_posix(),
        )
    )


def verify_ledger(checks: list[Check]) -> None:
    rows = read_csv_rows(LEDGER_PATH)
    checks.append(
        Check(
            f"LED-{len(checks)+1:03d}",
            "evidence_ledger",
            "formal evidence rows exist",
            "WARN" if len(rows) == 0 else "PASS",
            f"formal_evidence_rows={len(rows)}; zero means extraction is done but metric verification has not been booked",
            LEDGER_PATH.relative_to(ROOT).as_posix(),
        )
    )


def write_summary(checks: list[Check], numeric_rows: list[dict[str, Any]], table_rows: list[dict[str, Any]]) -> None:
    status_counts: dict[str, int] = {}
    for check in checks:
        status_counts[check.status] = status_counts.get(check.status, 0) + 1

    top_table = sorted(table_rows, key=lambda r: int(r["table_like_score"]), reverse=True)[:12]
    high_numeric = sorted(numeric_rows, key=lambda r: int(r["number_count"]), reverse=True)[:12]

    lines = [
        "# 多方法抽取核验报告",
        "",
        f"生成时间：{datetime.now().isoformat(timespec='seconds')}",
        "",
        "## 总结",
        "",
        f"- 检查项总数：{len(checks)}",
        f"- 状态统计：{status_counts}",
        "- 结论：原始文件、抽取 JSON、资料画像、关键词索引之间已完成多方法交叉检查。",
        "- 重要说明：`evidence_ledger.csv` 仍是正式指标入账表，目前没有指标行；这不是抽取失败，而是下一步待做。",
        "",
        "## 使用的方法",
        "",
        "1. 原始文件清单与 `data_catalog.csv` 对照。",
        "2. PDF 原始页数与抽取 JSON 页数对照。",
        "3. PPT 原始页数与抽取 JSON 页数对照。",
        "4. 原始文件重新抽文本，与保存的 JSON 文本长度逐页/逐幻灯片对照。",
        "5. 从保存 JSON 复算关键词出现次数，对照 `keyword_hits.csv`。",
        "6. PyMuPDF 原生表格检测 `find_tables()`。",
        "7. 数值密度和疑似表格页启发式扫描。",
        "8. PPT 原生表格和图片对象统计。",
        "9. 证据表入账状态检查。",
        "",
        "## 主要输出文件",
        "",
        "- `40_quality_evidence/verification/integrity_checks.csv`",
        "- `40_quality_evidence/verification/keyword_recount.csv`",
        "- `40_quality_evidence/verification/numeric_density.csv`",
        "- `40_quality_evidence/verification/table_candidates.csv`",
        "",
        "## 疑似表格/结构化数据优先复核页",
        "",
        "| 来源 | 单元 | 序号 | 表格分数 | 数字数 | 百分比数 | 金额数 | 片段 |",
        "|---|---:|---:|---:|---:|---:|---:|---|",
    ]
    for row in top_table:
        lines.append(
            f"| {row['source_file']} | {row['unit_type']} | {row['unit_index']} | {row['table_like_score']} | {row['numeric_count']} | {row['percent_count']} | {row['currency_count']} | {str(row['snippet']).replace('|', '/')} |"
        )

    lines.extend([
        "",
        "## 数值密度最高的页面/幻灯片",
        "",
        "| 来源 | 单元 | 序号 | 文本长度 | 数字数 | 百分比数 | 金额数 | 日期数 |",
        "|---|---:|---:|---:|---:|---:|---:|---:|",
    ])
    for row in high_numeric:
        lines.append(
            f"| {row['source_file']} | {row['unit_type']} | {row['unit_index']} | {row['text_length']} | {row['number_count']} | {row['percent_count']} | {row['currency_count']} | {row['date_count']} |"
        )

    lines.extend([
        "",
        "## 仍不能宣称完成的部分",
        "",
        "- 还没有逐条把关键指标写入 `evidence_ledger.csv`。",
        "- 还没有对 PPT 方案页的收益和供需缺口结论逐条回查强证据。",
        "- 还没有调用高德 API 做现实 POI 和路径核验。",
        "- 还没有对所有候选表格做人工抽样确认。",
    ])

    (VERIFY_DIR / "multi_method_verification.md").write_text("\n".join(lines), encoding="utf-8")


def main() -> None:
    VERIFY_DIR.mkdir(parents=True, exist_ok=True)
    checks: list[Check] = []
    numeric_rows: list[dict[str, Any]] = []
    table_rows: list[dict[str, Any]] = []

    check_catalog(checks)
    for path in source_files():
        if path.suffix.lower() == ".pdf":
            verify_pdf(path, checks, numeric_rows, table_rows)
        elif path.suffix.lower() in {".ppt", ".pptx"}:
            verify_ppt(path, checks, numeric_rows, table_rows)

    verify_keyword_hits(checks)
    verify_profile(checks)
    verify_ledger(checks)

    write_csv(
        VERIFY_DIR / "integrity_checks.csv",
        [check.__dict__ for check in checks],
        ["check_id", "scope", "method", "status", "detail", "evidence_path"],
    )
    write_csv(
        VERIFY_DIR / "numeric_density.csv",
        numeric_rows,
        [
            "source_file",
            "unit_type",
            "unit_index",
            "text_length",
            "table_like_score",
            "number_count",
            "percent_count",
            "currency_count",
            "date_count",
            "zero_like_count",
        ],
    )
    write_csv(
        VERIFY_DIR / "table_candidates.csv",
        table_rows,
        [
            "source_file",
            "unit_type",
            "unit_index",
            "pymupdf_table_count",
            "table_like_score",
            "numeric_count",
            "percent_count",
            "currency_count",
            "snippet",
        ],
    )
    write_summary(checks, numeric_rows, table_rows)
    print(f"checks={len(checks)} numeric_units={len(numeric_rows)} table_candidates={len(table_rows)} -> {VERIFY_DIR}")


if __name__ == "__main__":
    main()
